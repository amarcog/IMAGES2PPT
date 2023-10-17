# -*- coding: utf-8 -*-
"""
Author: Andrés Marco Giménez
This is a shiny app to generate automated ppt from indexed lined images with "_suffix".
"""
from shiny import App, Inputs, Outputs, Session, ui
from pptx import Presentation
from pptx.util import Cm
import numpy as np
import re
import tempfile
from pathlib import Path

app_ui = ui.page_fluid(

    ui.layout_sidebar(
        ui.panel_sidebar(
            {"style": "background-color: rgba(30, 144, 255, 0.4); font-weight: bold;"},
            ui.input_text("extension", "Extension of your images", placeholder=".tif .jpg or .png"),
            ui.input_numeric("img_size", "Size of your images in the ppt (cm)", value=6),
            ui.input_numeric("space_images", "Desired space between linked images (cm)", value=0.2),
            ui.input_numeric("margen_size", "Size of margens in slides (cm)", value=3),
            ui.input_numeric("n_images", "Number of images by slide", value=3),
            ui.input_text("channels", "Enter the name of your IDs separated by a SINGLE white space", placeholder="Ex. blue red green grays merge"),
            ui.input_file("IMAGES", "Choose your images", multiple=True),
            ui.download_button("downloadData", "Download Powerpoint", style="background: rgba(255, 148, 112, 0.4)"),
             
        
        ),

        ui.panel_main(
            ui.img(src="LOGO.png", style="width: 500px;"),
            ui.markdown(
            """
            Originary, IMAGES2PPT was developed to automatically order 
            images from microscopy in powerpoint presentations. However its usage 
            could be applied to any kind of images that have been labeled with a rule that 
            is compatible with the app workflow.

            Usually, microscopists capture images at different levels of magnification, 
            using various color channels (as in fluorescence imaging), 
            or creating replicas of their samples for further analysis. 
            These approaches ultimately generates a lot of linked images
            which is convenient to visualize side by side for comprehensive analysis.

            IMAGES2PPT allows to automatically control this kind of side by side visualizations
            in powerpoint format. **The only requirement is to build this labeling strategy:**

            # Samplename_ID.ext

            * **Samplename:** name of the sample that has been acquired at different magnifications, colors, replicas etc. It should be identical to all linked images.
            
            * **Underescore symbol \'_':** This is the key element for the app to consider every character encompassed between the last \'_' and the file's extension, as the magnification/channel identifier.

            * **Identifier 'ID':** characters encompassed between \_ and the file extension (.ext) that identify the different channels (i.e. blue, red, DAPI, ACTIN etc), magnification (i.e. 5x, 10x, 20x etc) or replica (Rep1, Rep2, Rep4...)

            * **File extension '.ext':** the extension of your image files (supported: .tif, .jpg, .png).

            ### Considering this labeling strategy, here there are useful examples for IMAGES2PPT usage:

            | Immunofluorescence   |      Magnifications  | Replicas |
            |----------------------|:--------------------:|------:|
            | Sample1_blue-DAPI.tif|  Sample1_5x.jpg         | Sample1_Rep1.png|
            | Sample1_green-ACTIN.tif |    Sample1_10x.jpg   |   Sample1_Rep2.png |
            | Sample1_red-WGA.tif | Sample1_20x.jpg          |    Sample1_Rep3.png |
            | Sample1_MERGE.tif | Sample1_30x.jpg            | Sample1_Rep4.png |
            | Sample2_blue-DAPI.tif|  Sample2_5x.jpg         | Sample2_Rep1.png|
            | Sample2_green-ACTIN.tif |    Sample2_10x.jpg   |   Sample2_Rep2.png |
            | Sample2_red-WGA.tif | Sample2_20x.jpg          |    Sample2_Rep3.png |
            | Sample2_MERGE.tif | Sample2_30x.jpg            | Sample2_Rep4.png |

            &copy; Copyright Andrés Marco Giménez
            """
        ),

        ),
    ),
)


def server(input: Inputs, output: Outputs, session: Session):

    @session.download(filename = 'Presentation.pptx')

    def downloadData():

        #Input data

        #(1)Set the name of your experiment
        #Inmunos_name = input.name()

        #(2)Set the size that you want for your images
        img_size = Cm(input.img_size())

        #(3)Set the size of the step that you want between images putted side by side
        space_images = Cm(input.space_images())

        #(4)Set the channels you want to plot
        channels_list = input.channels().split(" ")

        #(5)Set the number of images that you want to plot by slide
        n_images = input.n_images()

        #(6)Set the size of the margen that you want in your ppt
        margen_size = Cm(input.margen_size())

        #(7)Set the extension of your images
        extension = input.extension()

        #(8)Get the file names with their associated virtual path and
        #create a function to get virtual paths whenever a file name is called
        
        #if input.IMAGES() is None:
            #return "Please upload your images"
        f: list[FileInfo] = input.IMAGES()
        names_list = []
        path_list = []
        for i in range(0,len(f)):
            names_list.append(f[i]['name'])
            path_list.append(f[i]['datapath'])
        names_paths = {'File_name': names_list, 'Virtual_Path': path_list}

        def get_virtual_path(known_file_name):
            # Find the corresponding virtual path
            paired_virtual_path = names_paths['Virtual_Path'][names_paths['File_name'].index(known_file_name)]
            return paired_virtual_path

        #Since this point don't modify the script

        #Data derived from imput

        #(1)Calculate the number of channels

        n_channels = len(channels_list) #Calculate the number of channels

        #(2)Extracting image names

        pattern = r'.*_{}'.format(channels_list[0]+extension)
        image_list_full = [item for item in names_list if re.search(pattern, item)]
        image_list = [image.replace('_'+channels_list[0]+extension,'') for image in image_list_full]

        #(3)Function to define ranges that will serve 
        #to define the images from the image_list that will be included in each slide

        def create_ranges(List):
            
            List_rangles = []
            
            for j in range(0,len(List)-1):
                
                List_rangles.append(range(List[j],List[j+1]))
            
            return List_rangles    
            
        number_list = np.arange(0,len(image_list),n_images).tolist()
        number_list.append(len(image_list))
        range_list = create_ranges(number_list)

        #(4)Calculate the highest dimension of your images to know the reference
        #image dimension to define steps in coordinates inside a slide

        #4(b)Create a function with conditional statements to resize images to the highest dimension:

        def resize(img_sample,img_size):

            if img_sample.height == img_sample.width:
                
                img_sample.height = img_size
                img_sample.width = img_size
            
            elif img_sample.height > img_sample.width:
                
                Conversion_factor = img_size / Cm(img_sample.height)
                final_width = Cm(img_sample.width) * Conversion_factor
                img_sample.height = img_size
                img_sample.width = int(final_width)
            
            elif img_sample.width > img_sample.height:
                
                Conversion_factor = img_size / Cm(img_sample.width)
                final_height = Cm(img_sample.height) * Conversion_factor
                img_sample.width = img_size
                img_sample.height = int(final_height)

        #4a create a temporal presentation with a sample image inside

        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as temp_file:

            prs_temp = Presentation() 
            prs_temp.slide_width = n_channels*img_size + (n_channels-1)*space_images + margen_size*2
            prs_temp.slide_height = n_images*img_size + n_images * Cm(1) + margen_size*2
            blank_slide_layout = prs_temp.slide_layouts[6]
            slide_temp = prs_temp.slides.add_slide(blank_slide_layout)
            img_sample=slide_temp.shapes.add_picture(get_virtual_path(image_list_full[0]),margen_size,margen_size)
            #return "Your ppt is done"
        
            #cheched until here
        
            #4(c)get dimensions and calculate y steps

            resize(img_sample,img_size)

            y_step = img_sample.height
            x_step = img_sample.width


            del prs_temp
            del slide_temp
            del img_sample


            #Create final presentation

            prs = Presentation() #The presentation is called prs
            prs.slide_width = n_channels*x_step + (n_channels-1)*space_images + margen_size*2 #Set the width of the slides
            prs.slide_height = n_images*y_step + n_images * Cm(1) + margen_size*2 #Set the Height of the slides
            blank_slide_layout = prs.slide_layouts[6] #Set the layout of white empty background for the slides

            #Code creating the individual slides

            for r in range_list:
                
                slide = prs.slides.add_slide(blank_slide_layout) 

                for i,n in zip(range(0,n_images),r):
                
                    #Code for adding text box
                    
                    xpos_title = margen_size
                    
                    ypos_title = margen_size + y_step * (i) + Cm(1) * i
                    
                    txBox = slide.shapes.add_textbox(xpos_title, ypos_title, Cm(1), Cm(1))
                    tf = txBox.text_frame
                    tf.text = image_list[n]
                
                
                    #Code for adding Images. 
                    
                    for j in range(0,len(channels_list)):
                        
                        path = get_virtual_path(image_list[n] +'_'+channels_list[j]+extension) #path to image
                        
                        if j == 0:
                        
                            x_img = margen_size
                        
                        else:
                            
                            x_img = margen_size + x_step*(j)
                        
                        img=slide.shapes.add_picture(path,x_img+space_images*j,ypos_title+Cm(1))
                        
                        #Conditional statements to resize images to the highest dimension
                        
                        resize(img,img_size)

            prs.save(temp_file.name)
               
        return temp_file.name


www_dir = Path(__file__).parent / "www"

app = App(app_ui, server, static_assets=www_dir)